from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Colour palette ─────────────────────────────────────────────────────────
NAVY    = RGBColor(0x1a, 0x1a, 0x2e)   # dark navy – backgrounds / headings
BLUE    = RGBColor(0x0f, 0x34, 0x60)   # mid blue
ACCENT  = RGBColor(0x16, 0xA0, 0x85)   # teal accent
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT   = RGBColor(0xF4, 0xF6, 0xF9)   # very light grey
DARK    = RGBColor(0x2C, 0x2C, 0x2C)
YELLOW  = RGBColor(0xF3, 0x9C, 0x12)
RED     = RGBColor(0xC0, 0x39, 0x2B)
GREEN   = RGBColor(0x27, 0xAE, 0x60)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

W = prs.slide_width
H = prs.slide_height

# ── Helpers ────────────────────────────────────────────────────────────────

def blank_slide():
    blank_layout = prs.slide_layouts[6]
    return prs.slides.add_slide(blank_layout)

def bg(slide, color=NAVY):
    from pptx.util import Emu
    shape = slide.shapes.add_shape(1, 0, 0, W, H)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def rect(slide, x, y, w, h, fill=BLUE, alpha=None):
    shape = slide.shapes.add_shape(1, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background()
    return shape

def txt(slide, text, x, y, w, h, size=18, bold=False, color=WHITE,
        align=PP_ALIGN.LEFT, wrap=True):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return tb

def heading(slide, title, subtitle=None):
    """Full-width navy header bar with title + optional subtitle."""
    rect(slide, 0, 0, W, Inches(1.3), fill=NAVY)
    txt(slide, title, Inches(0.4), Inches(0.1), W - Inches(0.8), Inches(0.75),
        size=28, bold=True, color=WHITE)
    if subtitle:
        txt(slide, subtitle, Inches(0.4), Inches(0.78), W - Inches(0.8), Inches(0.42),
            size=14, color=RGBColor(0xAA, 0xCC, 0xEE))

def pill(slide, text, x, y, w, h, fill=ACCENT, tsize=11, tcolor=WHITE):
    shape = slide.shapes.add_shape(1, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.size = Pt(tsize)
    run.font.bold = True
    run.font.color.rgb = tcolor
    return shape

def card(slide, title, body_lines, x, y, w, h,
         title_fill=BLUE, body_fill=LIGHT, title_color=WHITE, body_color=DARK):
    """A titled card with bullet-ish content."""
    # title bar
    th = Inches(0.38)
    rect(slide, x, y, w, th, fill=title_fill)
    txt(slide, title, x + Inches(0.1), y + Inches(0.03), w - Inches(0.2), th,
        size=12, bold=True, color=title_color)
    # body
    rect(slide, x, y + th, w, h - th, fill=body_fill)
    body_text = "\n".join(body_lines)
    tb = slide.shapes.add_textbox(x + Inches(0.1), y + th + Inches(0.06),
                                   w - Inches(0.2), h - th - Inches(0.1))
    tf = tb.text_frame
    tf.word_wrap = True
    first = True
    for line in body_lines:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(2)
        run = p.add_run()
        run.text = line
        run.font.size = Pt(10.5)
        run.font.color.rgb = body_color

def mono(slide, lines, x, y, w, h):
    """Code-style box."""
    rect(slide, x, y, w, h, fill=RGBColor(0x1E, 0x1E, 0x2E))
    tb = slide.shapes.add_textbox(x + Inches(0.12), y + Inches(0.1),
                                   w - Inches(0.24), h - Inches(0.2))
    tf = tb.text_frame
    tf.word_wrap = False
    first = True
    for line in lines:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        run = p.add_run()
        run.text = line
        run.font.size = Pt(9.5)
        run.font.color.rgb = RGBColor(0xA8, 0xFF, 0x78)
        run.font.name = "Courier New"

def speaker_note(slide, text):
    slide.notes_slide.notes_text_frame.text = text


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE
# ══════════════════════════════════════════════════════════════════════════════
s = blank_slide()
bg(s, NAVY)

# large accent stripe at bottom
rect(s, 0, H - Inches(1.2), W, Inches(1.2), fill=ACCENT)

# main title
txt(s, "Hotel Management Backend",
    Inches(1), Inches(1.2), W - Inches(2), Inches(1.2),
    size=44, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
txt(s, "Databases for Developers — Oral Exam Presentation",
    Inches(1), Inches(2.5), W - Inches(2), Inches(0.7),
    size=20, color=RGBColor(0xAA, 0xCC, 0xFF), align=PP_ALIGN.CENTER)

# three db pills
for i, (label, col) in enumerate([
    ("MySQL  (Port 3307)", RGBColor(0x00, 0x75, 0x8F)),
    ("MongoDB  (Port 27017)", GREEN),
    ("Neo4j  (Port 7687)", RGBColor(0x80, 0x8A, 0xFF)),
]):
    pill(s, label, Inches(1.5 + i * 3.5), Inches(3.6), Inches(3.0), Inches(0.55),
         fill=col, tsize=14)

txt(s, "Spring Boot 4  ·  JWT Auth  ·  Docker Compose",
    Inches(1), Inches(4.5), W - Inches(2), Inches(0.5),
    size=14, color=RGBColor(0xFF, 0xFF, 0xFF), align=PP_ALIGN.CENTER)

speaker_note(s,
"TALEKORT — Slide 1: Intro\n"
"Præsenter kort projektet: Vi har bygget et Hotel Management System som backend-API i Spring Boot 4 med tre databaser – MySQL til strukturerede relationsdata, MongoDB til dokumenter og AI-logs, og Neo4j til grafrelationer. "
"Alt kører i Docker Compose og er sikret med JWT-baseret autentificering.")


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — SYSTEM ARCHITECTURE  (clean, no rotated shapes)
# ══════════════════════════════════════════════════════════════════════════════
s = blank_slide()
bg(s, LIGHT)
heading(s, "1 · System Architecture", "High-level overview — data flow and responsibilities")

# ── Spring Boot centre box ──────────────────────────────────────────────────
rect(s, Inches(3.8), Inches(1.55), Inches(5.2), Inches(3.6), fill=NAVY)
txt(s, "Spring Boot 4  ·  API :8080",
    Inches(3.9), Inches(1.65), Inches(5.0), Inches(0.5),
    size=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
# divider line inside API box
rect(s, Inches(3.9), Inches(2.22), Inches(4.8), Inches(0.03),
     fill=RGBColor(0x44, 0x55, 0x77))
for i, line in enumerate([
    "JWT Authentication  ·  BCrypt passwords",
    "Role-based access: ADMIN / STAFF / CLEANER",
    "@Transactional  (JPA + Neo4j managers)",
    "REST controllers  →  Service layer  →  DB",
]):
    txt(s, line, Inches(3.95), Inches(2.32) + i * Inches(0.62),
        Inches(4.9), Inches(0.55),
        size=10.5, color=RGBColor(0xBB, 0xCC, 0xEE), align=PP_ALIGN.CENTER)

# ── Client box (left) ───────────────────────────────────────────────────────
rect(s, Inches(0.2), Inches(2.6), Inches(2.8), Inches(1.5), fill=RGBColor(0x33, 0x33, 0x55))
txt(s, "Client\n(Browser / Frontend)",
    Inches(0.25), Inches(2.7), Inches(2.7), Inches(1.3),
    size=13, color=WHITE, align=PP_ALIGN.CENTER)
# connector line: client → API
rect(s, Inches(3.0), Inches(3.28), Inches(0.8), Inches(0.04),
     fill=RGBColor(0x77, 0x77, 0x99))
txt(s, "REST →", Inches(2.95), Inches(3.08), Inches(0.9), Inches(0.25),
    size=9, bold=True, color=RGBColor(0x55, 0x55, 0x88), align=PP_ALIGN.CENTER)

# ── AI Service box (bottom-left) ────────────────────────────────────────────
rect(s, Inches(0.2), Inches(4.55), Inches(2.8), Inches(1.2), fill=RGBColor(0x6A, 0x25, 0x7A))
txt(s, "AI Service\n(Ollama / LLM)",
    Inches(0.25), Inches(4.65), Inches(2.7), Inches(1.0),
    size=13, color=WHITE, align=PP_ALIGN.CENTER)
rect(s, Inches(3.0), Inches(4.97), Inches(0.8), Inches(0.04),
     fill=RGBColor(0x99, 0x55, 0xAA))
txt(s, "← HTTP →", Inches(2.93), Inches(4.77), Inches(0.95), Inches(0.25),
    size=9, bold=True, color=RGBColor(0x88, 0x44, 0x99), align=PP_ALIGN.CENTER)

# ── Three database boxes (right column) ─────────────────────────────────────
DB = [
    ("MySQL  :3307", "Source of truth · ACID · FK\nTransactions · JOINs · Triggers",
     RGBColor(0x00, 0x78, 0x96)),
    ("MongoDB  :27017", "Embedded guest docs · AI logs\nSchema-free · aggregation pipeline",
     RGBColor(0x1A, 0x8A, 0x45)),
    ("Neo4j  :7687", "Graph: nodes, relationships\nCypher · multi-hop traversal",
     RGBColor(0x50, 0x5A, 0xD0)),
]
db_x = Inches(9.35)
db_w = Inches(3.8)
db_h = Inches(1.08)
db_gap = Inches(0.22)
db_y0 = Inches(1.58)

for i, (title, body, col) in enumerate(DB):
    dy = db_y0 + i * (db_h + db_gap)
    # colour accent bar
    rect(s, db_x, dy, Inches(0.18), db_h, fill=col)
    # white body
    rect(s, db_x + Inches(0.18), dy, db_w - Inches(0.18), db_h,
         fill=WHITE)
    txt(s, title, db_x + Inches(0.26), dy + Inches(0.1),
        db_w - Inches(0.35), Inches(0.38),
        size=12, bold=True, color=DARK)
    txt(s, body, db_x + Inches(0.26), dy + Inches(0.5),
        db_w - Inches(0.35), Inches(0.52),
        size=9.5, color=RGBColor(0x44, 0x44, 0x44))
    # connector from API box right edge
    rect(s, Inches(9.0), dy + db_h/2 - Inches(0.02), Inches(0.35), Inches(0.04),
         fill=RGBColor(0x88, 0x88, 0xAA))
    txt(s, "→", Inches(8.95), dy + db_h/2 - Inches(0.18),
        Inches(0.35), Inches(0.35),
        size=14, bold=True, color=col, align=PP_ALIGN.CENTER)

# ── Migration banner ────────────────────────────────────────────────────────
rect(s, Inches(0.2), Inches(6.1), Inches(12.9), Inches(0.72),
     fill=RGBColor(0xE2, 0xEA, 0xF8))
txt(s, "POST /api/migrate  →  Copies data from MySQL into MongoDB & Neo4j   ·   "
       "Production: CDC (Change Data Capture) + Kafka for real-time, event-driven sync",
    Inches(0.4), Inches(6.2), Inches(12.5), Inches(0.55),
    size=10.5, color=DARK, align=PP_ALIGN.CENTER)

speaker_note(s,
"TALEKORT — Slide 2: Arkitektur\n"
"Systemet har én indgang: Spring Boot API'et på port 8080. Klienten sender REST-kald, API'et håndterer JWT-auth og delegerer til den rette database.\n"
"MySQL er kilden til sandheden – alle ændringer sker her. Via POST /api/migrate kopieres data til MongoDB og Neo4j.\n"
"AI-servicen (Ollama) er en separat container, og LLM-samtale-logs gemmes i MongoDB fordi de er fuldstændigt ustrukturerede.\n"
"Vi har to transaction managers: JpaTransactionManager (@Primary) til MySQL/JPA og Neo4jTransactionManager til Neo4j – Spring Boot 4 auto-konfigurerer ikke begge i multi-store setup.")


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — ERD  (proper entity boxes + relationship lines + cardinality)
# ══════════════════════════════════════════════════════════════════════════════
s = blank_slide()
bg(s, LIGHT)
heading(s, "2 · Relational Database Design — ERD",
        "Entities, PK/FK, cardinalities · 13 tables · InnoDB (ACID)")

# ── ERD helpers ─────────────────────────────────────────────────────────────
ENTITY_COL   = RGBColor(0x0A, 0x5A, 0x8A)
JUNCTION_COL = RGBColor(0x5A, 0x3A, 0x8A)
LINE_COL     = RGBColor(0x44, 0x44, 0x55)
CARD_COL     = ACCENT
SNAP_COL     = YELLOW

def erd_entity(sl, name, fields, ex, ey, ew, eh=None, hcol=ENTITY_COL):
    """Draw table box: dark header + light body with field list."""
    hh = Inches(0.31)
    n_fields = len(fields)
    if eh is None:
        eh = hh + Inches(0.06) + n_fields * Inches(0.175)
    rect(sl, Inches(ex), Inches(ey), Inches(ew), hh, fill=hcol)
    txt(sl, name, Inches(ex+0.06), Inches(ey+0.03),
        Inches(ew-0.12), hh - Inches(0.06),
        size=9, bold=True, color=WHITE)
    rect(sl, Inches(ex), Inches(ey)+hh, Inches(ew), Inches(eh)-hh,
         fill=RGBColor(0xE6, 0xF2, 0xFF))
    for fi, fld in enumerate(fields):
        txt(sl, fld,
            Inches(ex+0.07), Inches(ey)+hh + Inches(0.04) + fi*Inches(0.175),
            Inches(ew-0.14), Inches(0.175),
            size=7.8, color=DARK)

def erd_junction(sl, name, fields, ex, ey, ew):
    """Draw junction table in a different colour."""
    erd_entity(sl, name, fields, ex, ey, ew, hcol=JUNCTION_COL)

def hline(sl, x1, x2, y, lh=0.022):
    """Thin horizontal rule."""
    rect(sl, Inches(x1), Inches(y)-Inches(lh/2),
         Inches(x2-x1), Inches(lh), fill=LINE_COL)

def vline(sl, x, y1, y2, lh=0.022):
    """Thin vertical rule."""
    rect(sl, Inches(x)-Inches(lh/2), Inches(y1),
         Inches(lh), Inches(y2-y1), fill=LINE_COL)

def card_h(sl, lbl, x, y, anchor="left"):
    """Cardinality label on horizontal line."""
    offset = 0.04 if anchor == "left" else -0.28
    txt(sl, lbl, Inches(x+offset), Inches(y-0.23), Inches(0.28), Inches(0.2),
        size=11, bold=True, color=CARD_COL)

def card_v(sl, lbl, x, y, anchor="top"):
    """Cardinality label on vertical line."""
    offset = 0.05 if anchor == "top" else -0.2
    txt(sl, lbl, Inches(x+0.05), Inches(y+offset), Inches(0.25), Inches(0.2),
        size=11, bold=True, color=CARD_COL)

# ── Entity positions (x, y, w) — all in float inches ────────────────────────
# TOP ROW: room_type | season_rate
RT_X, RT_Y, RT_W = 5.55, 1.42, 2.2
SR_X, SR_Y, SR_W = 8.35, 1.42, 2.5

# MAIN ROW (y≈3.05): guest | reservation | room | room_cleaning_task | cleaner
G_X,  G_Y,  G_W  = 0.15, 3.05, 2.1
RV_X, RV_Y, RV_W = 2.65, 3.05, 2.75
RM_X, RM_Y, RM_W = RT_X, 3.05, RT_W   # aligned under room_type
CT_X, CT_Y, CT_W = SR_X, 3.05, 2.55   # aligned under season_rate
CL_X, CL_Y, CL_W = 11.2, 3.05, 2.0

# BOTTOM ROW: bill_item | bill  (aligned under guest/reservation)
BI_X, BI_Y, BI_W = 0.15, 5.05, 2.1
BL_X, BL_Y, BL_W = RV_X, 5.05, 2.3   # same x as reservation

# Entity heights (header 0.31 + body)
G_H  = 0.31 + 0.06 + 3*0.175   # 3 fields
RV_H = 0.31 + 0.06 + 4*0.175   # 4 fields
RM_H = 0.31 + 0.06 + 3*0.175
RT_H = 0.31 + 0.06 + 3*0.175
SR_H = 0.31 + 0.06 + 3*0.175
CT_H = 0.31 + 0.06 + 3*0.175
CL_H = 0.31 + 0.06 + 3*0.175
BL_H = 0.31 + 0.06 + 3*0.175
BI_H = 0.31 + 0.06 + 3*0.175

# Derived centres / edges
G_CY   = G_Y  + G_H/2;   G_R  = G_X  + G_W;   G_BOT = G_Y  + G_H
RV_CX  = RV_X + RV_W/2;  RV_CY = RV_Y + RV_H/2
RV_L   = RV_X;            RV_R  = RV_X + RV_W
RV_BOT = RV_Y + RV_H
RM_CX  = RM_X + RM_W/2;  RM_L  = RM_X;   RM_T  = RM_Y;   RM_TOP = RM_Y
RT_CX  = RT_X + RT_W/2;  RT_R  = RT_X + RT_W; RT_BOT = RT_Y + RT_H
SR_L   = SR_X
CT_L   = CT_X;            CT_R  = CT_X + CT_W
CL_L   = CL_X;            CT_CY = CT_Y + CT_H/2
BL_CX  = BL_X + BL_W/2;  BL_T  = BL_Y;   BL_L  = BL_X
BI_R   = BI_X + BI_W;    BI_CY = BI_Y + BI_H/2

# ── Draw entities ────────────────────────────────────────────────────────────
erd_entity(s, "guest", [
    "PK  guest_id  BIGINT  AUTO_INC",
    "first_name, last_name",
    "email UNIQUE, phone",
], G_X, G_Y, G_W)

erd_entity(s, "reservation", [
    "PK  reservation_id",
    "FK  guest_id, room_id",
    "FK  room_type_id, booked_rate_id",
    "status ENUM, booked_nightly_price ★",
], RV_X, RV_Y, RV_W)

erd_entity(s, "room", [
    "PK  room_id  INT",
    "FK  room_type_id",
    "room_number, status, occupied",
], RM_X, RM_Y, RM_W)

erd_entity(s, "room_type", [
    "PK  room_type_id  INT",
    "name  (Single/Double/Suite)",
    "max_occupancy",
], RT_X, RT_Y, RT_W, hcol=RGBColor(0x1A, 0x6A, 0x5A))

erd_entity(s, "season_rate", [
    "PK  rate_id  INT",
    "FK  room_type_id",
    "season, price_per_night",
], SR_X, SR_Y, SR_W, hcol=RGBColor(0x1A, 0x6A, 0x5A))

erd_entity(s, "room_cleaning_task", [
    "PK  task_id  INT",
    "FK  room_id",
    "task_type, scheduled_date",
], CT_X, CT_Y, CT_W, hcol=RGBColor(0x1A, 0x5A, 0x1A))

erd_entity(s, "cleaner", [
    "PK  cleaner_id  INT",
    "first_name, last_name",
    "phone, is_active",
], CL_X, CL_Y, CL_W, hcol=RGBColor(0x1A, 0x5A, 0x1A))

erd_entity(s, "bill", [
    "PK  bill_id  INT",
    "FK  reservation_id  UNIQUE",
    "total_amount, opened_at",
], BL_X, BL_Y, BL_W, hcol=RGBColor(0x7A, 0x3A, 0x00))

erd_entity(s, "bill_item", [
    "PK  bill_item_id  INT",
    "FK  bill_id",
    "item_type, unit_price, qty",
], BI_X, BI_Y, BI_W, hcol=RGBColor(0x7A, 0x3A, 0x00))

# Junction tables
erd_junction(s, "reservation_guest", [
    "PK  (reservation_id, guest_id)",
    "is_primary  BIT",
], 1.1, 4.35, 2.0)

erd_junction(s, "room_cleaning_assignment", [
    "PK  (task_id, cleaner_id)",
], 10.05, 4.35, 2.7)

# ── Relationship lines ───────────────────────────────────────────────────────
# 1. room_type ─1:N─ season_rate  (horizontal)
hline(s, RT_R, SR_L, RT_Y + RT_H/2)
card_h(s, "1", RT_R, RT_Y + RT_H/2, anchor="left")
card_h(s, "N", SR_L, RT_Y + RT_H/2, anchor="right")

# 2. room_type ─1:N─ room  (vertical, aligned on center x)
vline(s, RT_CX, RT_BOT, RM_TOP)
card_v(s, "1", RT_CX, RT_BOT, anchor="top")
card_v(s, "N", RT_CX, RM_TOP - 0.22, anchor="top")

# 3. season_rate ─1:N─ room_cleaning_task  (vertical, via room_type FK on season_rate)
#    Actually season_rate.room_type_id FK → room_type.  room_cleaning_task.room_id FK → room.
#    These don't connect directly. Draw room ─1:N─ room_cleaning_task instead:
hline(s, RM_X + RM_W, CT_L, RM_Y + RM_H/2)
card_h(s, "1", RM_X + RM_W, RM_Y + RM_H/2, anchor="left")
card_h(s, "N", CT_L, RM_Y + RM_H/2, anchor="right")

# 4. guest ─M:N─ reservation  (via reservation_guest)
#    Line from guest right → junction left, and junction right → reservation left
hline(s, G_R, 1.1, G_CY)                  # guest → junction left side
card_h(s, "M", G_R, G_CY, anchor="left")
hline(s, 1.1+2.0, RV_L, G_CY)             # junction right → reservation
card_h(s, "N", RV_L, G_CY, anchor="right")
# vertical drop from G_CY down to junction row
vline(s, 1.1 + 1.0, G_BOT, 4.35)         # centre of junction box vertically

# 5. reservation ─N:1─ room  (FK reservation.room_id)
hline(s, RV_R, RM_L, RV_Y + RV_H/2)
card_h(s, "N", RV_R, RV_Y + RV_H/2, anchor="left")
card_h(s, "1", RM_L, RV_Y + RV_H/2, anchor="right")

# 6. room_cleaning_task ─M:N─ cleaner  (via room_cleaning_assignment)
hline(s, CT_R, 10.05, CT_CY)
card_h(s, "M", CT_R, CT_CY, anchor="left")
hline(s, 10.05+2.7, CL_L, CT_CY)
card_h(s, "N", CL_L, CT_CY, anchor="right")
vline(s, 10.05 + 1.35, CT_Y + CT_H, 4.35)

# 7. reservation ─1:1─ bill  (vertical)
vline(s, RV_CX, RV_BOT, BL_T)
card_v(s, "1", RV_CX, RV_BOT, anchor="top")
card_v(s, "1", RV_CX, BL_T - 0.22, anchor="top")

# 8. bill ─1:N─ bill_item  (horizontal)
hline(s, BI_R, BL_L, BI_CY)
card_h(s, "N", BI_R, BI_CY, anchor="right")
card_h(s, "1", BL_L, BI_CY, anchor="right")

# ── Legend / annotation bar ──────────────────────────────────────────────────
rect(s, Inches(0.15), Inches(6.6), Inches(13.0), Inches(0.62),
     fill=RGBColor(0xE0, 0xE8, 0xF4))
legend_items = [
    (ENTITY_COL,   "Core entity"),
    (RGBColor(0x1A,0x6A,0x5A), "Type / rate"),
    (RGBColor(0x1A,0x5A,0x1A), "Housekeeping"),
    (RGBColor(0x7A,0x3A,0x00), "Billing"),
    (JUNCTION_COL, "Junction (M:N)"),
]
for li, (lc, ll) in enumerate(legend_items):
    lx = Inches(0.4 + li * 2.55)
    rect(s, lx, Inches(6.82), Inches(0.2), Inches(0.2), fill=lc)
    txt(s, ll, lx + Inches(0.26), Inches(6.77), Inches(2.1), Inches(0.28),
        size=9.5, color=DARK)

txt(s, "★ booked_nightly_price = historical SNAPSHOT (deliberate, not a 3NF violation)",
    Inches(0.4), Inches(6.62), Inches(12.5), Inches(0.22),
    size=8.5, color=RGBColor(0x80, 0x60, 0x00))

speaker_note(s,
"TALEKORT — Slide 3: ERD\n"
"Vi har 13 tabeller. De vigtigste pointer:\n"
"1. reservation_guest er en JUNCTION-tabel (lilla) til M:N relationen: én reservation kan have flere gæster, og én gæst kan optræde i flere reservationer. Den har COMPOSITE PK.\n"
"2. bill har UNIQUE FK til reservation – det er 1:1. Én reservation = præcis én regning.\n"
"3. booked_nightly_price gemmes som SNAPSHOT på reservationen. Selvom season_rate ændres i fremtiden, ved vi altid hvad gæsten faktisk betalte. Det er en bevidst designbeslutning – ikke en normaliseringsovertrædelse.\n"
"4. room_cleaning_assignment er endnu en junction-tabel: mange cleaners kan assignes til mange tasks.\n"
"5. Alle FK'er er konfigureret med ON DELETE CASCADE eller SET NULL for referentiel integritet.\n"
"6. room_type er en lookup-tabel (normalisering 3NF): name og max_occupancy hører til room_type, ikke i room eller reservation.")


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — NORMALIZATION
# ══════════════════════════════════════════════════════════════════════════════
s = blank_slide()
bg(s, LIGHT)
heading(s, "3 · Normalization — 1NF → 3NF",
        "How we reached our schema and key design choices")

# three columns
for i, (nf, color, title, bullets) in enumerate([
    ("1NF", RGBColor(0x00, 0x75, 0x8F), "Atomic values only",
     ["No lists/arrays in one column",
      "Each field = one indivisible value",
      "✓ credit_card_last4: only 4 digits",
      "✓ status: one ENUM value per row",
      "✓ Reservations in own table,\n   not comma-list in guest row"]),
    ("2NF", GREEN, "Full key dependency",
     ["All non-key fields depend on\nthe WHOLE primary key",
      "Relevant for COMPOSITE PKs",
      "✓ reservation_guest.is_primary\n   depends on BOTH columns",
      "✗ Would be wrong: guest_email\n   depends only on guest_id\n   → moved to guest table"]),
    ("3NF", ACCENT, "No transitive dependencies",
     ["Non-key fields must not depend\non other non-key fields",
      "A → B → C means C moves\nto its own table",
      "✓ room_type_name & max_occupancy\n   live in room_type, NOT in room\n   or reservation",
      "⚡ Exception: booked_nightly_price\n   is intentional snapshot,\n   not a violation"]),
]):
    cx = Inches(0.4 + i * 4.3)
    # header circle
    pill(s, nf, cx, Inches(1.4), Inches(1.2), Inches(0.65),
         fill=color, tsize=22)
    card(s, title, bullets, cx, Inches(2.15), Inches(4.0), Inches(4.5),
         title_fill=color, body_fill=WHITE, body_color=DARK)

speaker_note(s,
"TALEKORT — Slide 4: Normalisering\n"
"1NF: Hvert felt er atomisk – ingen lister. For eksempel gemmer vi credit_card_last4 som netop 4 cifre, ikke hele kortnummeret. Reservationer er i sin egen tabel, ikke som en tekstliste i guest-tabellen.\n"
"2NF: Gælder composite PKs. I reservation_guest afhænger is_primary af BEGGE nøgle-kolonner. Hvis vi havde lagt guest_email her, ville det kun afhænge af guest_id – det ville bryde 2NF.\n"
"3NF: Ingen transitive afhængigheder. room_type_name og max_occupancy hører til i room_type-tabellen, ikke i reservation. \n"
"Vigtig undtagelse: booked_nightly_price ser ud som transitiv afhængighed, men det er et bevidst historisk snapshot – prisen på bookingstidspunktet skal bevares selv om season_rate ændres.")


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — MONGODB
# ══════════════════════════════════════════════════════════════════════════════
s = blank_slide()
bg(s, LIGHT)
heading(s, "4 · MongoDB Design",
        "Collections, embedding vs. referencing, document model")

# left: document model sketch
card(s, "guest collection  (embedded)",
     [
         '{',
         '  "guestId": 5,',
         '  "firstName": "Magnus",',
         '  "email": "...",',
         '  "reservations": [          ← EMBEDDED',
         '    {',
         '      "reservationId": 42,',
         '      "checkInDate": "2024-06-01",',
         '      "roomType": "Double",',
         '      "bills": [             ← EMBEDDED',
         '        {',
         '          "totalAmount": 2400.00,',
         '          "items": [...]',
         '        }',
         '      ]',
         '    }',
         '  ]',
         '}',
     ],
     Inches(0.3), Inches(1.4), Inches(5.5), Inches(5.7),
     title_fill=GREEN, body_fill=RGBColor(0x1E, 0x1E, 0x2E), body_color=RGBColor(0xA8, 0xFF, 0x78))

# right: two decision cards
card(s, "EMBEDDED — Guest + Reservations + Bills",
     ["One document = full guest history",
      "Fetch everything in a single query",
      "No JOINs needed",
      "Best when: data is always read together",
      "Trade-off: updates must sync with MySQL"],
     Inches(6.2), Inches(1.4), Inches(6.8), Inches(2.3),
     title_fill=GREEN, body_fill=WHITE, body_color=DARK)

card(s, "REFERENCED — Rooms / Cleaners / Season Rates",
     ["Exist independently of reservations",
      "Shared across many documents",
      "Updated frequently (room status)",
      "Best when: high reuse, independent lifecycle",
      "Stored as separate collections with IDs"],
     Inches(6.2), Inches(3.9), Inches(6.8), Inches(2.0),
     title_fill=RGBColor(0x27, 0x6A, 0x8F), body_fill=WHITE, body_color=DARK)

card(s, "ai_interactions  (MongoDB only)",
     ["Unstructured LLM chat logs",
      "Schema-free — no fixed columns",
      "sessionId, userMessage, aiResponse,\ntokensUsed, model, timestamp"],
     Inches(6.2), Inches(6.0), Inches(6.8), Inches(0.85),
     title_fill=RGBColor(0x80, 0x30, 0x80), body_fill=WHITE, body_color=DARK)

speaker_note(s,
"TALEKORT — Slide 5: MongoDB\n"
"Det centrale designvalg i MongoDB er: skal relaterede data embeddes eller refereres?\n"
"Vi embedder reservationer og bills ind i guest-dokumentet fordi vi næsten altid henter dem sammen – det giver os én database-kald i stedet for multiple joins.\n"
"Rooms og cleaners er separate collections med referencer, fordi de eksisterer uafhængigt og deles af mange reservationer. Et rum opdateres hyppigt (status), og vi vil ikke opdatere det i hundredvis af embedded dokumenter.\n"
"ai_interactions gemmes kun i MongoDB fordi LLM-samtaler er fuldstændigt ustrukturerede – de passer ikke ind i et relationelt schema.")


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — NEO4J  (clean 3-row graph diagram)
# ══════════════════════════════════════════════════════════════════════════════
s = blank_slide()
bg(s, LIGHT)
heading(s, "5 · Neo4j Graph Model",
        "Nodes, relationships, properties — the hotel graph")

# ── Drawing helpers ───────────────────────────────────────────────────────────
N_HDR  = RGBColor(0x4A, 0x54, 0xCC)
N_BODY_C = RGBColor(0xEA, 0xEC, 0xFF)
REL_C  = RGBColor(0xAA, 0x55, 0x00)
LINE_C = RGBColor(0x55, 0x55, 0x88)

def nn(sl, lbl, props, x, y, w=2.1):
    """Neo4j node: dark header + light body."""
    hh = Inches(0.30)
    bh = Inches(0.06) + len(props) * Inches(0.168)
    rect(sl, Inches(x), Inches(y), Inches(w), hh, fill=N_HDR)
    txt(sl, lbl, Inches(x+0.06), Inches(y+0.025),
        Inches(w-0.12), hh - Inches(0.04), size=9, bold=True, color=WHITE)
    rect(sl, Inches(x), Inches(y)+hh, Inches(w), bh, fill=N_BODY_C)
    for i, p in enumerate(props):
        txt(sl, p, Inches(x+0.07), Inches(y)+hh+Inches(0.03)+i*Inches(0.168),
            Inches(w-0.14), Inches(0.168), size=7.5, color=DARK)

def rh(sl, x1, x2, y, name):
    """Horizontal relationship line (→ right) with label tag on line."""
    lw = Inches(0.024)
    rect(sl, Inches(x1), Inches(y) - lw/2, Inches(x2-x1), lw, fill=LINE_C)
    # arrowhead at right end
    txt(sl, "▶", Inches(x2 - 0.15), Inches(y - 0.13),
        Inches(0.17), Inches(0.17), size=7, color=LINE_C)
    # label pill centred on the line
    lw_lbl = 1.72
    lx = (x1 + x2) / 2 - lw_lbl / 2
    rect(sl, Inches(lx), Inches(y) - Inches(0.14), Inches(lw_lbl), Inches(0.26), fill=REL_C)
    txt(sl, name, Inches(lx), Inches(y) - Inches(0.145),
        Inches(lw_lbl), Inches(0.26),
        size=7.5, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

def rv(sl, x, y1, y2, name, up=False):
    """Vertical relationship line with label pill on right."""
    lw = Inches(0.024)
    rect(sl, Inches(x) - lw/2, Inches(y1), lw, Inches(y2 - y1), fill=LINE_C)
    # arrowhead
    if up:
        txt(sl, "▲", Inches(x - 0.09), Inches(y1 + 0.02),
            Inches(0.18), Inches(0.18), size=7, color=LINE_C)
    else:
        txt(sl, "▼", Inches(x - 0.09), Inches(y2 - 0.20),
            Inches(0.18), Inches(0.18), size=7, color=LINE_C)
    # label pill on the right
    mid_y = (y1 + y2) / 2
    lbl_h = 0.25
    rect(sl, Inches(x + 0.1), Inches(mid_y - lbl_h/2),
         Inches(2.15), Inches(lbl_h), fill=REL_C)
    txt(sl, name, Inches(x + 0.1), Inches(mid_y - lbl_h/2),
        Inches(2.15), Inches(lbl_h),
        size=7.5, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# ── Geometry constants ────────────────────────────────────────────────────────
# Node heights (header 0.30 + padding 0.06 + N×0.168)
def nh(n_props): return 0.30 + 0.06 + n_props * 0.168

# Row Y positions (top of first node in each row)
R1Y = 1.56
R2Y = 3.40
R3Y = 5.08

# Column X positions and widths
# Col 1: Guest      (x=0.20, w=2.00, cx=1.20)
# Col 2: Res/Bill/BillItem (x=2.68, w=2.50/2.25/2.25, cx≈3.93)
# Col 3: Room/Task/Cleaner (x=5.68, w=2.30, cx=6.83)
# Col 4: RoomType/SeasonRate (x=8.58, w=2.55, cx=9.855)

G_X, G_W   = 0.20,  2.00
RV_X, RV_W = 2.68,  2.50
RM_X, RM_W = 5.68,  2.30
RT_X, RT_W = 8.58,  2.55

BL_X, BL_W = 2.68,  2.30
CT_X, CT_W = 5.55,  2.60   # slightly wider for long name
SR_X, SR_W = 8.58,  2.55

BI_X, BI_W = 2.68,  2.30
CL_X, CL_W = 5.68,  2.30

# Column centres
G_CX  = G_X  + G_W/2    # 1.20
RV_CX = RV_X + RV_W/2   # 3.93
RM_CX = RM_X + RM_W/2   # 6.83
RT_CX = RT_X + RT_W/2   # 9.855
BL_CX = BL_X + BL_W/2   # 3.835 (≈ RV_CX close enough)
CT_CX = CT_X + CT_W/2   # 6.85
SR_CX = SR_CX = SR_X + SR_W/2  # 9.855
BI_CX = BI_X + BI_W/2   # 3.835
CL_CX = CL_X + CL_W/2   # 6.83

# Row 1 vertical centres (for horizontal connector y)
RH1_Y = R1Y + nh(3) / 2   # Guest/Res/Room have 3 props → centre ≈ 1.982

# Bottoms of nodes (for vertical connector start)
G_BOT  = R1Y + nh(3)       # 2.404
RV_BOT = R1Y + nh(3)       # 2.404
RM_BOT = R1Y + nh(3)       # 2.404
RT_BOT = R1Y + nh(2)       # 2.236

BL_BOT = R2Y + nh(2)       # 3.736
CT_BOT = R2Y + nh(2)       # 3.736
SR_BOT = R2Y + nh(2)       # 3.736

# ── Draw nodes ────────────────────────────────────────────────────────────────
# Row 1
nn(s, ":Guest",
   ["guestId  |  firstName", "lastName  |  email", "phone"],
   G_X, R1Y, G_W)

nn(s, ":Reservation",
   ["reservationId  |  status", "checkIn  |  checkOut", "nights  |  bookedPrice"],
   RV_X, R1Y, RV_W)

nn(s, ":Room",
   ["roomId  |  roomNumber", "roomStatus  |  cleanStatus", "occupied"],
   RM_X, R1Y, RM_W)

nn(s, ":RoomType",
   ["name  (Single/Double/Suite)", "maxOccupancy"],
   RT_X, R1Y, RT_W)

# Row 2
nn(s, ":Bill",
   ["billId  |  totalAmount", "openedAt  |  closedAt"],
   BL_X, R2Y, BL_W)

nn(s, ":RoomCleaningTask",
   ["taskId  |  taskType", "scheduledDate"],
   CT_X, R2Y, CT_W)

nn(s, ":SeasonRate",
   ["season  |  pricePerNight", "validFrom  |  validTo"],
   SR_X, R2Y, SR_W)

# Row 3
nn(s, ":BillItem",
   ["billItemId  |  itemType", "quantity  |  unitPrice", "lineTotal"],
   BI_X, R3Y, BI_W)

nn(s, ":Cleaner",
   ["cleanerId  |  firstName", "lastName  |  phone", "isActive"],
   CL_X, R3Y, CL_W)

# ── Draw relationship lines ───────────────────────────────────────────────────
# Horizontal row 1  (left edge of right node → right edge of left node)
rh(s, G_X+G_W,   RV_X,   RH1_Y, "[:BOOKED_BY]")
rh(s, RV_X+RV_W, RM_X,   RH1_Y, "[:ASSIGNED_TO]")
rh(s, RM_X+RM_W, RT_X,   RH1_Y, "[:IS_TYPE]")

# Vertical col 2: Reservation ↓ Bill ↓ BillItem
rv(s, RV_CX, RV_BOT, R2Y,   "[:HAS_BILL]")
rv(s, BL_CX, BL_BOT, R3Y,   "[:CONTAINS_ITEM]")

# Vertical col 3: Room ↓ RoomCleaningTask, Cleaner ↑ RoomCleaningTask
rv(s, RM_CX, RM_BOT, R2Y,   "[:HAS_CLEANING_TASK]")
rv(s, CL_CX, CT_BOT, R3Y,   "[:ASSIGNED_TO]", up=True)

# Vertical col 4: RoomType ↓ SeasonRate
rv(s, RT_CX, RT_BOT, R2Y,   "[:HAS_RATE]")

# ── Cypher example ────────────────────────────────────────────────────────────
txt(s, "Cypher — traverse 3 hops in one expression, no JOINs:",
    Inches(0.2), Inches(6.12), Inches(11.0), Inches(0.22),
    size=9, color=DARK)
mono(s, [
    "MATCH (g:Guest)-[:BOOKED_BY]->(res:Reservation)-[:ASSIGNED_TO]->(r:Room)",
    "WHERE g.guestId = 5",
    "RETURN g.firstName, res.status, r.roomNumber",
], Inches(0.2), Inches(6.34), Inches(11.0), Inches(0.88))

speaker_note(s,
"TALEKORT — Slide 6: Neo4j\n"
"Neo4j gemmer data som noder og relationer. Hvert node har et label fx :Guest eller :Room, og properties.\n"
"Pile-relationerne er navngivne: BOOKED_BY, ASSIGNED_TO, HAS_BILL osv. Man kan altid se retningen (pilen).\n"
"Fordelen: Cypher MATCH-udtryk følger pile naturligt. Fra Guest til Reservation til Room til Task til Cleaner – ét udtryk, ingen JOINs.\n"
"Det er specielt stærkt ved multi-hop spørgsmål: 'hvilke cleaners rengør rum booket af gæster med 3+ reservationer?' – trivielt i Cypher, 5 JOINs i SQL.")


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — CROSS-DB QUERY
# ══════════════════════════════════════════════════════════════════════════════
s = blank_slide()
bg(s, LIGHT)
heading(s, "6 · Cross-Database Query Comparison",
        "Business scenario: find all CONFIRMED reservations for room 101")

col_w = Inches(4.1)
cx = [Inches(0.2), Inches(4.5), Inches(8.8)]

titles  = ["MySQL  (SQL)", "MongoDB  (Aggregation)", "Neo4j  (Cypher)"]
colors_ = [RGBColor(0x00, 0x75, 0x8F), GREEN, RGBColor(0x60, 0x6A, 0xDE)]

sql_lines = [
    "SELECT g.first_name, g.last_name,",
    "       r.reference_no,",
    "       r.check_in_date",
    "FROM reservation r",
    "JOIN guest g",
    "  ON r.guest_id = g.guest_id",
    "JOIN room rm",
    "  ON r.room_id = rm.room_id",
    "WHERE rm.room_number = '101'",
    "  AND r.status = 'CONFIRMED';",
]
mongo_lines = [
    "db.guests.aggregate([",
    "  { $unwind: '$reservations' },",
    "  { $match: {",
    "    'reservations.room': '101',",
    "    'reservations.status':'CONFIRMED'",
    "  }},",
    "  { $project: {",
    "    firstName:1, lastName:1,",
    "    'reservations.referenceNo':1",
    "  }}",
    "])",
]
cypher_lines = [
    "MATCH (g:Guest)",
    "  -[:BOOKED_BY]->",
    "  (res:Reservation)",
    "  -[:ASSIGNED_TO]->",
    "  (r:Room)",
    "WHERE r.roomNumber = '101'",
    "  AND res.status = 'CONFIRMED'",
    "RETURN g.firstName,",
    "       g.lastName,",
    "       res.referenceNo",
]

for i, (lines, title, col) in enumerate(zip(
        [sql_lines, mongo_lines, cypher_lines], titles, colors_)):
    rect(s, cx[i], Inches(1.35), col_w, Inches(0.38), fill=col)
    txt(s, title, cx[i] + Inches(0.1), Inches(1.38), col_w - Inches(0.2), Inches(0.35),
        size=12, bold=True, color=WHITE)
    mono(s, lines, cx[i], Inches(1.73), col_w, Inches(3.5))

# comparison table
card(s, "Comparison",
     ["MySQL — explicit JOINs, set-based, ACID-strong, schema-enforced",
      "MongoDB — pipeline ($unwind → $match → $project), one document fetch, flexible schema",
      "Neo4j — pattern-match traversal (MATCH), no JOINs, best for multi-hop graph queries"],
     Inches(0.2), Inches(5.4), W - Inches(0.4), Inches(1.65),
     title_fill=NAVY, body_fill=WHITE, body_color=DARK)

speaker_note(s,
"TALEKORT — Slide 7: Cross-database query\n"
"Scenario: find alle bekræftede reservationer til rum 101.\n"
"MySQL: To JOINs – reservation til guest og til room. Klassisk SQL med WHERE-filter.\n"
"MongoDB: Aggregation pipeline med $unwind (folder embedded reservations-array ud), $match (filter) og $project (vælg felter). Alt i ét dokument-kald.\n"
"Neo4j: Cypher pattern MATCH følger pilene direkte. Ingen JOINs – vi traverserer grafen.\n"
"Pointen: alle tre giver samme svar, men med vidt forskellige styrker. SQL er bedst til komplekse joins over mange tabeller, MongoDB er bedst til embedded dokument-hierarkier, Neo4j er bedst til multi-hop relationsnavigation.")


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — TRANSACTIONS
# ══════════════════════════════════════════════════════════════════════════════
s = blank_slide()
bg(s, LIGHT)
heading(s, "7 · Transactions — ACID Guarantees",
        "Concrete example: Guest check-in — 3 operations, one transaction")

# ACID row
for i, (letter, name, desc, col) in enumerate([
    ("A", "Atomicity",   "All 3 steps succeed\nor none at all",      RGBColor(0x00, 0x75, 0x8F)),
    ("C", "Consistency", "FK constraints ensure\nno orphan records",  GREEN),
    ("I", "Isolation",   "Two guests booking\nsame room: only 1 wins",ACCENT),
    ("D", "Durability",  "Committed data\nsurvives server crash",     RED),
]):
    cx2 = Inches(0.3 + i * 3.25)
    pill(s, letter, cx2, Inches(1.4), Inches(0.6), Inches(0.7), fill=col, tsize=28)
    txt(s, name, cx2, Inches(2.2), Inches(3.1), Inches(0.35), size=12, bold=True, color=DARK)
    txt(s, desc, cx2, Inches(2.55), Inches(3.1), Inches(0.6), size=10, color=DARK)

# check-in transaction code
mono(s, [
    "BEGIN;",
    "",
    "  UPDATE reservation",
    "  SET status = 'CHECKED_IN'",
    "  WHERE reservation_id = 42;          -- step 1",
    "",
    "  UPDATE room SET occupied = 1,",
    "  room_status = 'OCCUPIED'",
    "  WHERE room_id = 15;                 -- step 2",
    "",
    "  INSERT INTO bill (reservation_id, opened_at, total_amount)",
    "  VALUES (42, NOW(), 0.00);           -- step 3",
    "",
    "COMMIT;   -- saves all three permanently",
    "-- If anything fails → ROLLBACK (nothing is saved)",
], Inches(0.3), Inches(3.3), Inches(7.5), Inches(3.8))

card(s, "Spring Boot — @Transactional",
     ["Service methods annotated with @Transactional",
      "Spring wraps the method in BEGIN...COMMIT",
      "Any exception triggers automatic ROLLBACK",
      "Neo4j uses separate Neo4jTransactionManager bean"],
     Inches(8.0), Inches(3.3), Inches(5.1), Inches(3.8),
     title_fill=NAVY, body_fill=WHITE, body_color=DARK)

speaker_note(s,
"TALEKORT — Slide 8: Transaktioner\n"
"En transaktion er en gruppe operationer der enten alle lykkes eller alle fejler – aldrig halvt.\n"
"Konkret eksempel: en gæst checker ind. Det kræver tre operationer: (1) reservation-status sættes til CHECKED_IN, (2) rummet markeres OCCUPIED, (3) en bill oprettes. Alle tre skal lykkes – ellers er systemet i inkonsistent tilstand.\n"
"ACID: A=alt-eller-intet, C=fra valid til valid tilstand, I=concurrent booking af samme rum – kun én lykkes, D=data forbliver gemt ved crash.\n"
"I Spring Boot bruger vi @Transactional på service-klassen. Vi har to transaction managers: JpaTransactionManager for MySQL/JPA og Neo4jTransactionManager for Neo4j, og JPA-manageren er @Primary.")


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — INDEXES
# ══════════════════════════════════════════════════════════════════════════════
s = blank_slide()
bg(s, LIGHT)
heading(s, "8 · Indexes & Performance",
        "Key indexes, design decisions, and what we left out")

# analogy
pill(s, "Index = book's back index — jump straight to the page, no full scan",
     Inches(0.3), Inches(1.42), Inches(12.7), Inches(0.45),
     fill=NAVY, tsize=12, tcolor=WHITE)

# index table
rows_idx = [
    ["Table",       "Index",                      "Type",      "Why?"],
    ["guest",       "idx_guest_email",             "B-tree",    "Login lookup — most frequent query"],
    ["guest",       "idx_guest_name (last, first)","Composite", "Name search — multi-column sort"],
    ["reservation", "idx_reservation_dates",       "Composite", "Date-range availability queries"],
    ["reservation", "idx_reservation_status",      "B-tree",    "Filter CONFIRMED / CHECKED_IN"],
    ["reservation", "idx_reservation_guest",       "B-tree",    "FK join: guest's reservations"],
    ["bill",        "idx_bill_reservation",        "B-tree",    "FK join performance"],
    ["room",        "idx_room_status",             "B-tree",    "Filter AVAILABLE rooms"],
    ["audit_log",   "idx_audit_table/op/time",     "B-tree",    "Fast audit trail queries"],
]

from pptx.util import Inches as I2
col_w2 = [I2(2.1), I2(3.6), I2(2.0), I2(5.0)]
yy = Inches(2.05)
header_drawn = False
for ri, row in enumerate(rows_idx):
    for ci, (cell, cw) in enumerate(zip(row, col_w2)):
        fill = NAVY if ri == 0 else (WHITE if ri % 2 == 0 else LIGHT)
        fc   = WHITE if ri == 0 else DARK
        bx   = Inches(0.3) + sum(col_w2[:ci])
        rect(s, bx, yy, cw, Inches(0.37), fill=fill)
        txt(s, cell, bx + Inches(0.05), yy + Inches(0.04),
            cw - Inches(0.1), Inches(0.3), size=9.5,
            bold=(ri == 0), color=fc)
    yy += Inches(0.37)

# not indexed note
card(s, "What we did NOT index — and why",
     ["bill_item.description — text search unused, low selectivity",
      "occupied BIT — only 2 values (low cardinality): full scan = ~50% of rows anyway",
      "Indexes slow down INSERT/UPDATE — only add where the SELECT gain outweighs the write cost"],
     Inches(0.3), Inches(5.7), Inches(12.7), Inches(1.5),
     title_fill=RED, body_fill=WHITE, body_color=DARK)

speaker_note(s,
"TALEKORT — Slide 9: Indexes\n"
"Et index er som en bogs indeks bagerst – i stedet for at læse alle rækker hopper databasen direkte til den rigtige. Uden index: full table scan, O(n). Med index: O(log n).\n"
"B-tree index er standardtypen og god til exact match og range queries.\n"
"Composite index dækker flere kolonner – vigtigt: man skal altid søge på den første kolonne i indexet for at det virker.\n"
"Vi valgte IKKE at indexere occupied-kolonnen fordi den kun har to værdier (lav cardinality) – databasen henter alligevel halvdelen af rækkerne, så indexet hjælper ikke. Det er en bevidst designbeslutning.")


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — SECURITY & AUDITING
# ══════════════════════════════════════════════════════════════════════════════
s = blank_slide()
bg(s, LIGHT)
heading(s, "9 · Security & Auditing",
        "JWT authentication, BCrypt, SQL injection prevention, audit triggers")

card(s, "JWT — JSON Web Token (stateless auth)",
     ["1. POST /api/auth/login → server returns signed token",
      "2. Client sends: Authorization: Bearer <token>",
      "3. Server verifies signature — NO database call needed",
      "Token contains: username, role, expiry  (HS512 signed)"],
     Inches(0.3), Inches(1.4), Inches(6.2), Inches(1.9),
     title_fill=RGBColor(0x00, 0x75, 0x8F), body_fill=WHITE, body_color=DARK)

card(s, "BCrypt — one-way password hashing",
     ['"admin123" → BCrypt → "$2a$10$N9qo8uLO..."',
      "Cannot be reversed — no plain-text passwords stored",
      "Salt built-in: same password → different hash each time"],
     Inches(6.7), Inches(1.4), Inches(6.3), Inches(1.9),
     title_fill=RGBColor(0x60, 0x20, 0x20), body_fill=WHITE, body_color=DARK)

card(s, "SQL Injection Prevention",
     ['UNSAFE: "SELECT * FROM guest WHERE email = \'" + input + "\'"',
      '  Attack: input = "\' OR \'1\'=\'1"  → returns ALL guests!',
      "SAFE: Spring Data JPA parameterized queries:",
      '  findByEmail(String email)',
      '  → SELECT * FROM guest WHERE email = ?',
      "  Input is treated as DATA, never as SQL code"],
     Inches(0.3), Inches(3.45), Inches(6.2), Inches(2.8),
     title_fill=RED, body_fill=WHITE, body_color=DARK)

card(s, "Roles & Access Control",
     ["ADMIN — full access (delete, create, view all)",
      "STAFF — can read/update reservations, not delete",
      "CLEANER — can only view/update cleaning tasks",
      "Enforced via Spring Security @PreAuthorize"],
     Inches(6.7), Inches(3.45), Inches(6.3), Inches(1.7),
     title_fill=NAVY, body_fill=WHITE, body_color=DARK)

card(s, "Database Audit Triggers (05_audit.sql)",
     ["Triggers on: reservation, bill, bill_item, guest, room",
      "Each INSERT/UPDATE/DELETE → row appended to audit_log",
      "Stores: old_values (JSON), new_values (JSON), changed_by, changed_at",
      "App does NOTHING — database fires trigger automatically"],
     Inches(6.7), Inches(5.3), Inches(6.3), Inches(1.9),
     title_fill=ACCENT, body_fill=WHITE, body_color=DARK)

mono(s, [
    "CREATE TRIGGER tr_audit_reservation_update",
    "AFTER UPDATE ON reservation FOR EACH ROW",
    "INSERT INTO audit_log(table_name, operation_type,",
    "  record_id, old_values, new_values, changed_by)",
    "VALUES ('reservation','UPDATE', NEW.reservation_id,",
    "  JSON_OBJECT('status', OLD.status),",
    "  JSON_OBJECT('status', NEW.status), CURRENT_USER());",
], Inches(0.3), Inches(3.45), Inches(6.0), Inches(2.8))

speaker_note(s,
"TALEKORT — Slide 9: Security og Auditing\n"
"JWT: stateless authentication. Brugeren logger ind og får et token. Fremtidige kald sender dette token i header – serveren verificerer signaturen uden at slå op i databasen. Tokenet indeholder username, rolle og udløbstid.\n"
"BCrypt: passwords gemmes aldrig i plain text. BCrypt er en one-way hash – kan ikke gå baglæns. Salt forhindrer rainbow table-angreb.\n"
"SQL injection: den klassiske sikkerhedsrisiko. Vi bruger Spring Data JPA parameterized queries – input behandles altid som data, aldrig som SQL-kode.\n"
"Audit triggers: en trigger er SQL der kører automatisk ved INSERT/UPDATE/DELETE. Applikationen behøver ikke aktivt at logge – databasen gør det. Hvert audit-entry gemmer hvad der stod FØR og EFTER ændringen som JSON.")


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — TRADE-OFFS & REFLECTIONS
# ══════════════════════════════════════════════════════════════════════════════
s = blank_slide()
bg(s, LIGHT)
heading(s, "10 · Trade-offs & Reflections",
        "Why three databases — and what we learned")

# 3 why-cards
for i, (db, col, reasons) in enumerate([
    ("Why MySQL?", RGBColor(0x00, 0x75, 0x8F),
     ["Structured, relational hotel data",
      "ACID guarantees for payments & bookings",
      "Complex multi-table queries (views, JOINs)",
      "Referential integrity via FK constraints"]),
    ("Why MongoDB?", GREEN,
     ["Full guest history in one document call",
      "Unstructured AI/LLM interaction logs",
      "Schema-free — add fields without migration",
      "Read-heavy workload with nested data"]),
    ("Why Neo4j?", RGBColor(0x60, 0x6A, 0xDE),
     ["Multi-hop relationship traversal",
      "Natural for 'who cleaned what booked by whom'",
      "Graph queries > recursive SQL for deep paths",
      "Potential for recommendation features"]),
]):
    card(s, db, reasons,
         Inches(0.3 + i * 4.35), Inches(1.4), Inches(4.1), Inches(2.5),
         title_fill=col, body_fill=WHITE, body_color=DARK)

# what we would do differently
card(s, "What we would do differently",
     ["Date types: stored dates as String in Neo4j instead of native date/datetime types → less efficient, cannot use Cypher date math",
      "Sync strategy: manual POST /api/migrate → in production would use Change Data Capture (CDC) with Kafka for automatic, event-driven propagation",
      "Multi-store Spring Boot: auto-configuration does NOT handle JPA + Neo4j transaction managers together → required explicit @Bean definitions for both (JpaTransactionManager @Primary + Neo4jTransactionManager named bean)",
      "Data consistency: MongoDB/Neo4j counts can drift from MySQL if migration is not re-run after every change"],
     Inches(0.3), Inches(4.05), W - Inches(0.6), Inches(3.1),
     title_fill=YELLOW, body_fill=WHITE, body_color=DARK)

speaker_note(s,
"TALEKORT — Slide 10: Trade-offs og refleksioner\n"
"MySQL er valgt fordi hoteldata er inherent relationelt – en gæst booker et rum til en specifik pris, og alt hænger logisk sammen. ACID er kritisk for finansielle transaktioner.\n"
"MongoDB er valgt for to formål: embedded dokument-hierarkier til read-optimized gæste-historik, og schema-free storage til AI-logs.\n"
"Neo4j er valgt for dens naturlige håndtering af forbindelser og traversals – spørgsmål om 'hvem er forbundet til hvad' er trivielle i Cypher men kræver mange JOINs i SQL.\n"
"Hvad vi ville gøre anderledes:\n"
"1. Native date-typer i Neo4j i stedet for strings.\n"
"2. CDC/Kafka til automatisk sync i stedet for manuel migration.\n"
"3. Vi opdagede undervejs at Spring Boot 4 ikke auto-konfigurerer Neo4jTransactionManager når JPA er til stede – det krævede eksplicit konfiguration af to separate transaction managers.")


# ── Save ───────────────────────────────────────────────────────────────────
import os
OUT = os.path.join(os.path.dirname(__file__), "Hotel_Eksamen_v2.pptx")
prs.save(OUT)
print(f"PowerPoint gemt: {OUT}")
